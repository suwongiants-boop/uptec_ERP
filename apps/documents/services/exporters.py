from __future__ import annotations

import io
import textwrap
from html import escape
from pathlib import Path

from django.http import FileResponse, Http404
from django.utils.text import slugify
from docx import Document as WordDocument
from docx.shared import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


class DocumentExportService:
    def build_response(self, document, export_format: str) -> FileResponse:
        preview_pages = (document.assembled_content or {}).get("preview_pages", [])
        if not preview_pages:
            raise Http404("No preview available for export.")

        export_format = export_format.lower()
        if export_format == "html":
            buffer, content_type, filename = self._export_html(document, preview_pages)
        elif export_format == "docx":
            buffer, content_type, filename = self._export_docx(document, preview_pages)
        elif export_format == "pptx":
            buffer, content_type, filename = self._export_pptx(document, preview_pages)
        elif export_format == "pdf":
            buffer, content_type, filename = self._export_pdf(document, preview_pages)
        else:
            raise Http404("Unsupported export format.")

        buffer.seek(0)
        return FileResponse(buffer, as_attachment=True, filename=filename, content_type=content_type)

    def _export_html(self, document, preview_pages):
        title = escape(document.title)
        html = self.render_html(document, preview_pages, title=title)
        return io.BytesIO(html.encode("utf-8")), "text/html; charset=utf-8", f"{self._slug(document.title)}.html"

    def render_html(self, document, preview_pages, title=None) -> str:
        title = title or escape(document.title)
        chunks = [
            "<!doctype html>",
            "<html lang='ko'>",
            "<head>",
            "<meta charset='utf-8'>",
            "<meta name='viewport' content='width=device-width, initial-scale=1'>",
            f"<title>{title}</title>",
            "<style>",
            "body{font-family:'Segoe UI','Apple SD Gothic Neo',sans-serif;background:#f4f7fb;color:#152033;margin:0;padding:40px;}",
            ".page{max-width:920px;margin:0 auto 24px;background:#fff;border:1px solid #dce4ef;border-radius:24px;padding:36px;box-shadow:0 12px 32px rgba(21,32,51,.08);}",
            ".page-label{display:inline-block;padding:6px 10px;border-radius:999px;background:#e9f2ff;color:#0057a4;font-weight:700;font-size:12px;}",
            "h1,h2{margin:14px 0 12px;}",
            "p{line-height:1.75;margin:0 0 12px;}",
            "ul{margin:14px 0 0 20px;line-height:1.75;}",
            ".image-note{margin-top:16px;padding:14px;border-radius:16px;background:#f7fbff;border:1px dashed #bdd1eb;color:#49566d;}",
            "</style>",
            "</head>",
            "<body>",
        ]
        for page in preview_pages:
            chunks.append("<section class='page'>")
            chunks.append(f"<span class='page-label'>{escape(page.get('page_label', ''))}</span>")
            chunks.append(f"<h1>{escape(page.get('title', ''))}</h1>")
            for paragraph in page.get("body", []):
                chunks.append(f"<p>{escape(paragraph)}</p>")
            bullets = page.get("bullets", [])
            if bullets:
                chunks.append("<ul>")
                for bullet in bullets:
                    chunks.append(f"<li>{escape(bullet)}</li>")
                chunks.append("</ul>")
            image_note = page.get("image_note", "")
            if image_note:
                chunks.append(f"<div class='image-note'>{escape(image_note)}</div>")
            chunks.append("</section>")
        chunks.extend(["</body>", "</html>"])
        return "\n".join(chunks)

    def _export_docx(self, document, preview_pages):
        output = io.BytesIO()
        word = WordDocument()
        word.core_properties.title = document.title
        title = word.add_heading(document.title, level=0)
        title.alignment = 1

        meta = (document.assembled_content or {}).get("preview_meta", {})
        word.add_paragraph(
            f"Language: {meta.get('language_label', '한국어')} / Pages: {meta.get('page_count', len(preview_pages))}"
        )

        for index, page in enumerate(preview_pages):
            word.add_heading(f"{page.get('page_label', '')} - {page.get('title', '')}", level=1)
            for paragraph in page.get("body", []):
                p = word.add_paragraph(paragraph)
                p.paragraph_format.space_after = Pt(10)
            for bullet in page.get("bullets", []):
                word.add_paragraph(bullet, style="List Bullet")
            if page.get("image_note"):
                note = word.add_paragraph()
                run = note.add_run(f"Image note: {page['image_note']}")
                run.italic = True
            if index != len(preview_pages) - 1:
                word.add_page_break()

        word.save(output)
        return (
            output,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            f"{self._slug(document.title)}.docx",
        )

    def _export_pptx(self, document, preview_pages):
        output = io.BytesIO()
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]

        for page in preview_pages:
            slide = presentation.slides.add_slide(blank_layout)
            title_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(0.4), PptInches(8.2), PptInches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = page.get("title", "")
            title_frame.paragraphs[0].font.size = PptPt(24)
            title_frame.paragraphs[0].font.bold = True

            subtitle_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(1.1), PptInches(2.0), PptInches(0.4))
            subtitle_box.text_frame.text = page.get("page_label", "")
            subtitle_box.text_frame.paragraphs[0].font.size = PptPt(11)

            body_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(1.6), PptInches(8.0), PptInches(4.8))
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            first = True
            for paragraph in page.get("body", []):
                current = body_frame.paragraphs[0] if first else body_frame.add_paragraph()
                current.text = paragraph
                current.font.size = PptPt(16)
                current.space_after = PptPt(10)
                first = False
            for bullet in page.get("bullets", []):
                current = body_frame.add_paragraph()
                current.text = bullet
                current.level = 1
                current.font.size = PptPt(14)

            note_box = slide.shapes.add_textbox(PptInches(0.6), PptInches(6.35), PptInches(8.3), PptInches(0.7))
            note_box.text_frame.text = page.get("image_note", "")
            note_box.text_frame.paragraphs[0].font.size = PptPt(11)

        presentation.save(output)
        return (
            output,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            f"{self._slug(document.title)}.pptx",
        )

    def _export_pdf(self, document, preview_pages):
        output = io.BytesIO()
        pages = []
        font_title = self._load_font(46)
        font_text = self._load_font(25)
        font_small = self._load_font(20)

        for page in preview_pages:
            canvas = Image.new("RGB", (1240, 1754), "#ffffff")
            draw = ImageDraw.Draw(canvas)
            draw.rounded_rectangle((46, 46, 1194, 1708), radius=28, outline="#d5dfec", width=3)

            draw.text((90, 82), page.get("page_label", ""), fill="#0057a4", font=font_small)
            draw.text((90, 140), page.get("title", ""), fill="#152033", font=font_title)

            y = 250
            for paragraph in page.get("body", []):
                for line in self._wrap_text(paragraph, 62):
                    draw.text((90, y), line, fill="#23324a", font=font_text)
                    y += 38
                y += 18

            if page.get("bullets"):
                y += 8
                for bullet in page["bullets"]:
                    for index, line in enumerate(self._wrap_text(bullet, 54)):
                        prefix = "• " if index == 0 else "  "
                        draw.text((108, y), prefix + line, fill="#23324a", font=font_text)
                        y += 36
                    y += 10

            if page.get("image_note"):
                y += 24
                draw.rounded_rectangle((90, y, 1140, min(y + 220, 1600)), radius=18, outline="#bfd2eb", width=2)
                for line in self._wrap_text("Image note: " + page["image_note"], 60):
                    draw.text((120, y + 24), line, fill="#49566d", font=font_small)
                    y += 30

            pages.append(canvas)

        if not pages:
            raise Http404("No preview pages available for PDF export.")

        pages[0].save(output, format="PDF", save_all=True, append_images=pages[1:])
        return output, "application/pdf", f"{self._slug(document.title)}.pdf"

    def _slug(self, title: str) -> str:
        return slugify(title, allow_unicode=True) or "document-export"

    def _load_font(self, size: int):
        font_candidates = [
            Path("C:/Windows/Fonts/malgun.ttf"),
            Path("C:/Windows/Fonts/segoeui.ttf"),
        ]
        for candidate in font_candidates:
            if candidate.exists():
                return ImageFont.truetype(str(candidate), size=size)
        return ImageFont.load_default()

    def _wrap_text(self, value: str, width: int) -> list[str]:
        cleaned = " ".join((value or "").split())
        if not cleaned:
            return []
        return textwrap.wrap(cleaned, width=width, break_long_words=False, replace_whitespace=False)
